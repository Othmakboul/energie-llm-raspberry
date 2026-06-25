"""PowerPoint point tuteurs — 25/06/2026 — 3 modèles complets."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BG    = RGBColor(0x0F,0x0F,0x1A); PANEL = RGBColor(0x16,0x16,0x2A)
RED   = RGBColor(0xE9,0x4F,0x37); BLUE  = RGBColor(0x00,0xB4,0xD8)
GREEN = RGBColor(0x2E,0xCC,0x71); GOLD  = RGBColor(0xF4,0xD0,0x3F)
PURP  = RGBColor(0x9B,0x59,0xB6); WHITE = RGBColor(0xFF,0xFF,0xFF)
GREY  = RGBColor(0x88,0x88,0x99)

W = Inches(13.33); H = Inches(7.5)
prs = Presentation(); prs.slide_width=W; prs.slide_height=H
blank = prs.slide_layouts[6]

def bg(s):
    r = s.shapes.add_shape(1,0,0,W,H)
    r.fill.solid(); r.fill.fore_color.rgb=BG; r.line.fill.background()

def rect(s,x,y,w,h,c):
    r=s.shapes.add_shape(1,x,y,w,h)
    r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background(); return r

def txt(s,t,x,y,w,h,sz=13,bold=False,color=WHITE,align=PP_ALIGN.LEFT,italic=False):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=t; r.font.size=Pt(sz)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=color; return tb

def header(s,title,sub=None):
    bg(s); rect(s,0,0,W,Inches(1.25),PANEL); rect(s,0,0,Inches(0.12),H,RED)
    txt(s,title,Inches(0.3),Inches(0.1),Inches(12.8),Inches(0.8),sz=30,bold=True)
    if sub: txt(s,sub,Inches(0.3),Inches(0.82),Inches(12.8),Inches(0.38),sz=13,color=BLUE)

def img(s,path,x,y,w,h):
    if os.path.exists(path): s.shapes.add_picture(path,x,y,w,h)

def num(s,n): txt(s,f"{n}/5",Inches(12.6),Inches(7.08),Inches(0.7),Inches(0.32),sz=11,color=GREY,align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Vue d'ensemble + résultats clés
# ══════════════════════════════════════════════════════════════
s=prs.slides.add_slide(blank)
header(s,"Point d'avancement — Semaine 3/6",
       "Abidi Amine & Othmane Makboul  •  LISTIC / USMB  •  25/06/2026")

# 4 campagnes
campaigns = [
    (RED,  "Llama-3.2-1B",  "19/06","3 quantif (Q3/Q4/Q8)","1 215 mesures"),
    (BLUE, "Qwen2.5-1.5B",  "24/06","3 quantif (Q3/Q4/Q8)","1 215 mesures"),
    (PURP, "Gemma-3-1B",    "25/06","3 quantif (Q3/Q4/Q8)","1 215 mesures"),
    (GOLD, "Prompts longs", "17/06","126 → 4 012 tokens input","18 mesures"),
]
for i,(col,nom,date,desc,tot) in enumerate(campaigns):
    cx=Inches(0.28)+i*Inches(3.22)
    rect(s,cx,Inches(1.42),Inches(3.08),Inches(3.3),PANEL)
    rect(s,cx,Inches(1.42),Inches(3.08),Inches(0.07),col)
    txt(s,nom,  cx+Inches(0.12),Inches(1.55),Inches(2.85),Inches(0.42),sz=14,bold=True,color=col)
    txt(s,date, cx+Inches(0.12),Inches(1.97),Inches(2.85),Inches(0.3), sz=11,color=GREY)
    txt(s,desc, cx+Inches(0.12),Inches(2.28),Inches(2.85),Inches(0.8), sz=12,color=WHITE)
    txt(s,tot,  cx+Inches(0.12),Inches(3.1), Inches(2.85),Inches(0.38),sz=14,bold=True,color=col)

# Résultats clés
rect(s,Inches(0.28),Inches(4.9),Inches(12.75),Inches(2.45),PANEL)
rect(s,Inches(0.28),Inches(4.9),Inches(12.75),Inches(0.06),GREEN)
txt(s,"Résultats confirmés — 3 645 mesures sur Pi 5",
    Inches(0.45),Inches(4.97),Inches(8.0),Inches(0.38),sz=14,bold=True,color=GREEN)

kpis=[
    (GREEN,"★ Composition optimale","Gemma-3-1B Q3_K_M  =  0,350 J/tok @ 15,1 tok/s  →  meilleure config toutes architectures"),
    (GOLD, "Loi linéaire",          "E = coût_fixe + α × tokens  —  r > 0,997 sur les 3 modèles et toutes quantifications"),
    (RED,  "Sweet spot ≠ selon archi","Llama/Qwen : Q4 < Q3 < Q8  |  Gemma : Q3 << Q8 ≈ Q4  (pattern inversé, Q3 gagne chez Gemma)"),
    (BLUE, "CodeCarbon",            "Surestime PMIC de +18 à +32 %  (RAPL absent sur ARM → estimation par TDP)"),
    (PURP, "Prompts longs",         "Effet négligeable < 500 tokens, mais ×16 entre 126 et 4 012 tokens input  (r=0,996)"),
]
for i,(col,label,val) in enumerate(kpis):
    cy=Inches(5.38)+i*Inches(0.37)
    txt(s,f"▶  {label} :",Inches(0.45),cy,Inches(2.5),Inches(0.34),sz=11,bold=True,color=col)
    txt(s,val,            Inches(2.95),cy,Inches(10.0),Inches(0.34),sz=11,color=WHITE)

num(s,1)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Comparaison 3 modèles Q4 + J/tok par quantif
# ══════════════════════════════════════════════════════════════
s=prs.slides.add_slide(blank)
header(s,"Résultat 1 — Comparaison 3 architectures × quantification",
       "Llama-3.2-1B  •  Qwen2.5-1.5B  •  Gemma-3-1B")

img(s,"data/fig1_3modeles_q4.png",Inches(0.2),Inches(1.35),Inches(12.9),Inches(3.0))
img(s,"data/fig2_jtok_par_quantif.png",Inches(0.2),Inches(4.45),Inches(12.9),Inches(2.85))

num(s,2)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Loi linéaire + carte efficacité
# ══════════════════════════════════════════════════════════════
s=prs.slides.add_slide(blank)
header(s,"Résultat 2 — Loi linéaire & carte d'efficacité",
       "E = coût_fixe + α × tokens  |  r > 0,997  |  Gemma Q3 = configuration optimale")

img(s,"data/fig3_lineaire_et_carte.png",Inches(0.2),Inches(1.35),Inches(12.9),Inches(4.9))

# Encart synthèse lois linéaires
rect(s,Inches(0.2),Inches(6.35),Inches(12.9),Inches(1.0),PANEL)
rows=[
    ("Llama Q4", "E = 0,8 + 0,448·tok", RED),
    ("Qwen  Q4", "E = 1,2 + 0,517·tok", BLUE),
    ("Gemma Q4", "E = 2,4 + 0,464·tok", PURP),
    ("Gemma Q3★","E = 0,4 + 0,347·tok", GREEN),
    ("Gemma Q8", "E = 0,4 + 0,483·tok", PURP),
    ("r > 0,997","sur toutes configs",   WHITE),
]
for i,(lbl,eq,col) in enumerate(rows):
    cx=Inches(0.35)+i*Inches(2.15)
    txt(s,lbl,cx,Inches(6.42),Inches(2.1),Inches(0.3),sz=10,bold=True,color=col)
    txt(s,eq, cx,Inches(6.73),Inches(2.1),Inches(0.5),sz=10,color=WHITE)

num(s,3)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Prompts longs + CodeCarbon
# ══════════════════════════════════════════════════════════════
s=prs.slides.add_slide(blank)
header(s,"Résultat 3 — Prompts longs & fiabilité CodeCarbon",
       "Effet prompt visible > 500 tokens  |  CodeCarbon surestime jusqu'à +32%")

img(s,"data/fig4_prompts_et_cc.png",Inches(0.2),Inches(1.35),Inches(12.9),Inches(4.75))

rect(s,Inches(0.2),Inches(6.2),Inches(12.9),Inches(1.15),PANEL)
obs=[
    (GOLD, "Prompts longs",
     "Négligeable < 500 tok  |  ×4 à 1000 tok  |  ×16 à 4000 tok  |  E = 2,1 + 0,152·input_tok  (r=0,996)"),
    (BLUE, "CodeCarbon vs PMIC",
     "Surestime de +18% (Llama Q3) à +32% (Qwen Q8)  —  Pas utilisable seul sur ARM sans PMIC de référence"),
]
for i,(col,lbl,val) in enumerate(obs):
    cy=Inches(6.27)+i*Inches(0.47)
    txt(s,f"▶ {lbl} :",Inches(0.4),cy,Inches(2.4),Inches(0.38),sz=11,bold=True,color=col)
    txt(s,val,          Inches(2.8),cy,Inches(10.2),Inches(0.38),sz=11,color=WHITE)

num(s,4)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Tableau final + suite + question
# ══════════════════════════════════════════════════════════════
s=prs.slides.add_slide(blank)
header(s,"Synthèse finale & Suite",
       "Composition optimale identifiée  •  S4 : analyse + interface + rapport")

# Tableau synthèse 3 modèles
rect(s,Inches(0.25),Inches(1.42),Inches(7.8),Inches(4.95),PANEL)
rect(s,Inches(0.25),Inches(1.42),Inches(7.8),Inches(0.06),WHITE)
txt(s,"Tableau synthèse — Q4_K_M (référence commune)",
    Inches(0.4),Inches(1.52),Inches(7.5),Inches(0.35),sz=12,bold=True,color=WHITE)

headers_t=["Modèle","Énergie (J)","J/token","tok/s","Rang"]
col_ws=[Inches(2.1),Inches(1.4),Inches(1.3),Inches(1.2),Inches(1.0)]
xstarts=[Inches(0.35)]+[Inches(0.35)+sum(col_ws[:i]) for i in range(1,5)]
for j,(h,xp) in enumerate(zip(headers_t,xstarts)):
    txt(s,h,xp,Inches(1.92),col_ws[j],Inches(0.3),sz=10,bold=True,color=GREY,align=PP_ALIGN.CENTER)

# Tableau trié par rang (1er → 3e)
table_data=[
    ("🥇  Llama-3.2-1B", RED,  "28,8 J","0,456","13,4","1er"),
    ("🥈  Gemma-3-1B",   PURP, "30,9 J","0,491","11,2","2e"),
    ("🥉  Qwen2.5-1.5B", BLUE, "33,3 J","0,532","11,0","3e"),
]
for i,(nom,col,e,jt,sp,rg) in enumerate(table_data):
    cy=Inches(2.26)+i*Inches(0.5)
    bg_c=RGBColor(0x0A,0x0A,0x18) if i%2==0 else PANEL
    rect(s,Inches(0.35),cy,Inches(12.6),Inches(0.46),bg_c)
    for j,(val,xp) in enumerate(zip([nom,e,jt,sp,rg],xstarts)):
        txt(s,val,xp,cy+Pt(5),col_ws[j],Inches(0.38),
            sz=12,bold=(j==0),color=col if j==0 else WHITE,align=PP_ALIGN.CENTER)

# Séparateur + config optimale globale
rect(s,Inches(0.35),Inches(3.78),Inches(12.6),Inches(0.06),GREEN)
rect(s,Inches(0.35),Inches(3.88),Inches(12.6),Inches(1.55),RGBColor(0x05,0x15,0x05))
txt(s,"★  Configuration optimale toutes architectures confondues :",
    Inches(0.5),Inches(3.95),Inches(7.0),Inches(0.38),sz=13,bold=True,color=GREEN)
txt(s,"Gemma-3-1B  Q3_K_M",
    Inches(7.5),Inches(3.95),Inches(5.0),Inches(0.38),sz=15,bold=True,color=PURP)
txt(s,"0,350 J/token  @  15,1 tok/s",
    Inches(0.5),Inches(4.38),Inches(6.0),Inches(0.45),sz=20,bold=True,color=GREEN)
txt(s,"−23% vs Llama Q4  |  −34% vs Qwen Q4  |  Pattern : Q3 << Q8 ≈ Q4",
    Inches(0.5),Inches(4.82),Inches(12.0),Inches(0.38),sz=12,color=GREY)

# Suite S4
rect(s,Inches(0.35),Inches(5.55),Inches(12.6),Inches(0.06),GOLD)
rect(s,Inches(0.35),Inches(5.65),Inches(12.6),Inches(1.7),PANEL)
txt(s,"Suite — Semaine 4",Inches(0.5),Inches(5.72),Inches(4.0),Inches(0.38),sz=13,bold=True,color=GOLD)
etapes=[
    ("Interface Streamlit", "Visualisation interactive de toutes les campagnes"),
    ("Rapport",             "Brouillon : introduction + méthodologie + résultats"),
    ("Option",              "Tester modèles 3B si temps disponible"),
]
for i,(titre,desc) in enumerate(etapes):
    cx=Inches(0.5)+i*Inches(4.25)
    txt(s,f"{i+1}. {titre}",cx,Inches(6.12),Inches(4.0),Inches(0.32),sz=12,bold=True,color=GOLD)
    txt(s,desc,             cx,Inches(6.45),Inches(4.0),Inches(0.7), sz=11,color=WHITE)

num(s,5)

out=r"C:\Users\Amine\Desktop\amiine_Rasbery_Pi\point_tuteurs_25juin.pptx"
prs.save(out); print(f"Sauvegardé : {out}")
